"""
PPT 识别接口 — PPT/图片/PDF → 文字
"""
import logging
import os
import uuid
from fastapi import APIRouter, UploadFile, File, HTTPException
from pydantic import BaseModel
from ai_services.ocr.service import OCRService
from app.models.database import Note, get_session

logger = logging.getLogger(__name__)

router = APIRouter()
_ocr_service = None

def get_ocr():
    """懒加载获取 OCR 服务"""
    global _ocr_service
    if _ocr_service is None:
        _ocr_service = OCRService()
    return _ocr_service

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp"}
PDF_EXTS = {".pdf"}
PPT_EXTS = {".ppt", ".pptx"}


# ==================== PPT/图片/PDF（瞬间返回 JSON）====================

@router.post("/ocr/demo")
async def ocr_demo(file: UploadFile = File(...)):
    ext = (os.path.splitext(file.filename or "file.png")[1]).lower()

    content = await file.read()

    if not content:
        return {"error": "上传的文件为空，请选择有效的文件"}

    if len(content) > 50 * 1024 * 1024:
        return {"error": "文件超过 50MB 限制，请压缩后再试"}

    save_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}{ext}")
    with open(save_path, "wb") as f:
        f.write(content)

    try:
        if ext in PPT_EXTS:
            text = _extract_pptx(save_path)
            if text:
                return {"type": "ppt", "text": text}
            pdf_path = _convert_ppt_to_pdf(save_path)
            if pdf_path and os.path.exists(pdf_path):
                text = await _extract_pdf_text(pdf_path)
                os.remove(pdf_path)
                if text:
                    return {"type": "ppt", "text": text}
            return {"type": "ppt", "text": "未从 PPT 中提取到文字内容"}

        elif ext in IMAGE_EXTS:
            try:
                page = await get_ocr().recognize_image(save_path)
                return {"type": "image", "text": page.full_text or "未识别到文字"}
            except Exception as ocr_err:
                logger.error(f"OCR识别异常: {ocr_err}")
                return {"error": f"OCR识别失败: {str(ocr_err)}"}

        elif ext in PDF_EXTS:
            try:
                text = await _extract_pdf_text(save_path)
                return {"type": "pdf", "text": text or "未识别到文字"}
            except Exception as pdf_err:
                logger.error(f"PDF识别异常: {pdf_err}")
                return {"error": f"PDF识别失败: {str(pdf_err)}"}

        else:
            return {"error": f"不支持: {ext}，支持 ppt/pptx/png/jpg/webp/bmp/pdf"}

    except Exception as e:
        logger.error(f"处理异常: {e}")
        return {"error": f"处理失败: {str(e)}"}

    finally:
        if os.path.exists(save_path):
            os.remove(save_path)


# ==================== OCR 转笔记 ====================

class OcrToNoteRequest(BaseModel):
    text: str
    ai_enhanced: bool = True


@router.post("/ocr/to-note")
async def ocr_to_note(data: OcrToNoteRequest):
    """将 OCR 识别结果创建为笔记（AI 增强）"""
    raw_text = data.text.strip()
    if not raw_text:
        raise HTTPException(status_code=400, detail="识别内容为空，无法创建笔记")

    lines = raw_text.split("\n")
    title = lines[0][:100] if lines else "PPT 识别笔记"
    tags = ["PPT"]

    md_content = raw_text
    ai_enhanced_applied = False
    ai_enhanced_message = ""

    if data.ai_enhanced:
        try:
            cfg_session = get_session()
            try:
                from app.models.database import UserConfig

                def _get_cfg(key: str) -> str | None:
                    c = cfg_session.query(UserConfig).filter(UserConfig.key == key).first()
                    return c.value if c else None

                api_key = _get_cfg("llm_api_key")
                base_url = _get_cfg("llm_base_url") or "https://api.deepseek.com/v1"
                model = _get_cfg("llm_model") or "deepseek-chat"
            finally:
                cfg_session.close()

            if api_key:
                from ai_services.note_generator.service import NoteGenerator
                generator = NoteGenerator()
                result = generator.enhance_via_openai(
                    api_key=api_key,
                    base_url=base_url,
                    model=model,
                    title=title,
                    content=raw_text,
                )
                md_content = result["content"]
                title = result["title"]
                ai_enhanced_applied = True
                logger.info(f"OCR 笔记 AI 增强完成")
            else:
                logger.warning("OCR 转笔记：未配置 API Key，跳过 AI 增强")
                ai_enhanced_message = "未配置 API Key，跳过 AI 增强"
        except Exception as e:
            logger.error(f"OCR 笔记 AI 增强失败: {e}", exc_info=True)
            md_content = raw_text
            title = lines[0][:100] if lines else "PPT 识别笔记"
            ai_enhanced_message = f"AI 增强失败: {str(e)}"

    session = get_session()
    try:
        note = Note(
            title=title,
            content=md_content,
            tags=tags,
            raw_markdown=md_content,
        )
        session.add(note)
        session.commit()
        session.refresh(note)
        return {
            "id": note.id,
            "title": note.title,
            "message": "笔记创建成功",
            "ai_enhanced_applied": ai_enhanced_applied,
            "ai_enhanced_message": ai_enhanced_message,
        }
    except Exception as e:
        session.rollback()
        raise HTTPException(status_code=500, detail=f"创建笔记失败: {str(e)}")
    finally:
        session.close()


# ==================== 工具函数 ====================

async def _extract_pdf_text(path: str) -> str:
    """使用 PyMuPDF 将 PDF 每页转图片，再用 RapidOCR 识别"""
    import fitz
    doc = fitz.open(path)
    ocr = get_ocr()
    pages = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(dpi=200)
        img_data = pix.tobytes("png")
        result = await ocr.recognize_image(img_data, is_bytes=True)
        text = result.full_text.strip()
        if text:
            pages.append(f"第 {page_num + 1} 页\n{text}")
        else:
            pages.append(f"第 {page_num + 1} 页\n（空白页或无文字）")
    doc.close()
    return "\n\n".join(pages)


def _extract_pptx(path: str) -> str:
    """使用 python-pptx 提取 PPT 幻灯片中的文字"""
    from pptx import Presentation
    try:
        prs = Presentation(path)
    except Exception:
        return ""
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"第 {i} 页"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append("  ".join(cells))
            if shape.shape_type == 13:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
        if slide.has_notes_slide:
            n = slide.notes_slide.notes_text_frame.text.strip()
            if n:
                parts.append(f"备注: {n}")
        slides.append("\n".join(parts))
    return "\n\n".join(slides)


def _convert_ppt_to_pdf(ppt_path: str) -> str | None:
    """将 PPT/PPTX 转为 PDF（需要 Windows + Microsoft Office）"""
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        base = os.path.splitext(ppt_path)[0]
        pdf_path = base + ".pdf"
        deck = powerpoint.Presentations.Open(ppt_path)
        deck.SaveAs(pdf_path, 32)
        deck.Close()
        powerpoint.Quit()
        return pdf_path
    except Exception as e:
        logger.warning(f"PPT→PDF 转换失败（可选功能，不影响文字提取）: {e}")
        return None